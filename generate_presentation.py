import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

project_dir = r"c:\Users\Personal\PycharmProjects\PythonProject"
os.chdir(project_dir)

files = {
    'Semana 26': 'Produccion Astroflores a la Semana 26_Seguimiento.xlsx',
    'Semana 27': 'Produccion Astroflores BL25-26-27-28 a la Semana 27_Seguimiento.xlsx',
    'Semana 28': 'Produccion Astroflores a la Semana 28_seguimiento.xlsx',
}

frames = {}
for label, file_name in files.items():
    path = os.path.join(project_dir, file_name)
    xls = pd.ExcelFile(path)
    sheet_name = xls.sheet_names[0]
    df = pd.read_excel(path, sheet_name=sheet_name)
    frames[label] = df

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Modelo de Proyección de Producción'
slide.placeholders[1].text = 'Astroflores\nSeguimiento Semanas 26, 27 y 28'

# Slide 2 - README summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Resumen del proyecto'
body = slide.shapes.placeholders[1].text_frame
body.clear()
body.word_wrap = True
text = (
    '• Proyección semanal por variedad (Bloque&Varid)\n'
    '• Uso de patrones históricos similares\n'
    '• Modelo RandomForestRegressor\n'
    '• Variables de ciclo y patrón: Producción lag 12, cambio vs lag 12 y semana de ciclo\n'
    '• Evaluación comparativa con datos reales de seguimiento'
)
body.paragraphs[0].text = text

# Slide 3 - metrics summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Evaluación con datos reales'
body = slide.shapes.placeholders[1].text_frame
body.clear()
body.word_wrap = True
p = body.paragraphs[0]
p.text = ''
for label, df in frames.items():
    diff_col = None
    for candidate in ['%Dif', '%DIF', 'Dif']:
        if candidate in df.columns:
            diff_col = candidate
            break
    if diff_col is None:
        continue
    diffs = pd.to_numeric(df[diff_col], errors='coerce').fillna(0)
    avg_diff = diffs.mean()
    abs_avg_diff = abs(diffs).mean()
    p = body.add_paragraph()
    p.text = f'{label}: registros={len(df)}, promedio %Dif={avg_diff:.2f}%, promedio absoluto %Dif={abs_avg_diff:.2f}%'
    p.level = 0

# Slide 4 - top varieties table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Top variedades por producción real - Semana 28'
df28 = frames['Semana 28']
summary = (
    df28.groupby('Bloque&Varid', as_index=False)['Produccion']
    .sum()
    .sort_values('Produccion', ascending=False)
    .head(8)
)
rows, cols = len(summary) + 1, 4
shape = slide.shapes.add_table(rows, cols, Inches(
    0.4), Inches(1.6), Inches(12.0), Inches(4.5))
table = shape.table
headers = ['Variedad', 'Producción real', 'Tallos/m2', 'Semana']
for col_idx, header in enumerate(headers):
    table.cell(0, col_idx).text = header

for idx, row in summary.reset_index(drop=True).iterrows():
    variedad = str(row['Bloque&Varid'])
    prod = float(row['Produccion'])
    tallos = float(df28.loc[df28['Bloque&Varid'] ==
                   variedad, 'Tallos/m2'].iloc[0])
    table.cell(idx + 1, 0).text = variedad
    table.cell(idx + 1, 1).text = f'{prod:,.0f}'
    table.cell(idx + 1, 2).text = f'{tallos:.2f}'
    table.cell(idx + 1, 3).text = '28'

# Slide 5 - key findings
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Hallazgos clave'
body = slide.shapes.placeholders[1].text_frame
body.clear()
body.word_wrap = True
text = (
    '• Los archivos de seguimiento muestran producción real semanal por variedad.\n'
    '• Se observa variabilidad fuerte entre variedades y semanas.\n'
    '• La proyección se apoya en historial, patrón y ciclo de 12 semanas.\n'
    '• El modelo se evalúa frente a datos reales para ajustar su desempeño.'
)
body.paragraphs[0].text = text

out_path = os.path.join(project_dir, 'presentacion_modelo_astroflores.pptx')
prs.save(out_path)
print(f'Presentación creada en: {out_path}')
